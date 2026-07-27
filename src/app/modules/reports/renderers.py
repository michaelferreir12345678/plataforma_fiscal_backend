"""Renderização de PDF, Excel e PPTX a partir de um documento fiscal auditável."""

from __future__ import annotations

import json
from io import BytesIO
from typing import Any

MIME_TYPES = {
    "pdf": "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _text(value: Any) -> str:
    if value is None:
        return "—"
    if isinstance(value, bool):
        return "Sim" if value else "Não"
    return str(value)


def _source_label(source_refs: list[dict[str, Any]]) -> str:
    return "; ".join(
        " · ".join(
            filter(
                None,
                (
                    str(ref.get("relatorio") or ""),
                    str(ref.get("anexo") or ""),
                    str(ref.get("periodo") or ""),
                    f"v{ref.get('versao_entrega')}" if ref.get("versao_entrega") else "",
                ),
            )
        )
        for ref in source_refs
    )


def render(document: dict[str, Any], formato: str) -> bytes:
    if formato == "pdf":
        return _render_pdf(document)
    if formato == "xlsx":
        return _render_xlsx(document)
    if formato == "pptx":
        return _render_pptx(document)
    raise ValueError(f"Formato não suportado: {formato}")


def _render_pdf(document: dict[str, Any]) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )

    buffer = BytesIO()
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="Institutional",
            parent=styles["Title"],
            textColor=colors.HexColor("#173F2A"),
            fontSize=17,
            leading=21,
            alignment=TA_CENTER,
        )
    )
    styles.add(
        ParagraphStyle(
            name="SmallAudit",
            parent=styles["BodyText"],
            fontSize=7.5,
            leading=10,
            textColor=colors.HexColor("#4F6257"),
        )
    )
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=15 * mm,
        rightMargin=15 * mm,
        topMargin=16 * mm,
        bottomMargin=16 * mm,
        title=document["titulo"],
        author=document["cabecalho"].get("responsavel") or "Plataforma Fiscal",
    )
    story: list[Any] = []
    header = document["cabecalho"]
    story.append(Paragraph(header["organizacao"], styles["Institutional"]))
    story.append(Paragraph(header["ente"], styles["Heading2"]))
    story.append(Paragraph(document["titulo"], styles["Heading1"]))
    story.append(
        Paragraph(
            f"Período: {document['periodo']} · consulta as_of: {document['as_of']} · "
            f"gerado em: {document['gerado_em']}",
            styles["SmallAudit"],
        )
    )
    story.append(Spacer(1, 8))

    if document["dados_incompletos"]:
        story.append(Paragraph("Dados incompletos ou defasados", styles["Heading2"]))
        issue_rows = [["Tipo", "Item", "Sinalização explícita"]]
        issue_rows.extend(
            [issue["tipo"], issue["codigo"], issue["mensagem"]]
            for issue in document["dados_incompletos"]
        )
        issue_table = Table(issue_rows, colWidths=[24 * mm, 40 * mm, 112 * mm], repeatRows=1)
        issue_table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#F4D8B5")),
                    ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#5B3214")),
                    ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D4B68C")),
                    ("FONT", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 7.5),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]
            )
        )
        story.append(issue_table)
        story.append(Spacer(1, 10))

    story.append(Paragraph("Indicadores fiscais", styles["Heading2"]))
    metric_rows = [["Indicador", "Valor", "Situação", "Fonte / período / versão", "as_of"]]
    for metric in document["metricas"]:
        metric_rows.append(
            [
                metric["rotulo"],
                metric["valor_formatado"],
                metric["status"],
                _source_label(metric["source_refs"]),
                metric.get("as_of") or "—",
            ]
        )
    metrics_table = Table(
        metric_rows,
        colWidths=[43 * mm, 27 * mm, 22 * mm, 57 * mm, 31 * mm],
        repeatRows=1,
    )
    metrics_table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#173F2A")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONT", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BCC8C0")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F7F5")]),
            ]
        )
    )
    story.append(metrics_table)

    if document.get("conformidade"):
        story.append(Spacer(1, 10))
        story.append(Paragraph("Entregas e conformidade", styles["Heading2"]))
        rows = [["Relatório", "Período", "Prazo", "Status", "Versão"]]
        rows.extend(
            [
                item["relatorio"],
                item["periodo"],
                item.get("prazo") or "—",
                item["status"],
                item.get("versao_entrega") or "—",
            ]
            for item in document["conformidade"]
        )
        table = Table(rows, colWidths=[35 * mm, 30 * mm, 35 * mm, 35 * mm, 35 * mm])
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E5EDE8")),
                    ("FONT", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.3, colors.HexColor("#BCC8C0")),
                ]
            )
        )
        story.append(table)

    story.append(PageBreak())
    story.append(Paragraph("Memória de cálculo e rastreabilidade", styles["Heading1"]))
    for metric in document["metricas"]:
        components = json.dumps(metric["memoria"].get("componentes", {}), ensure_ascii=False)
        story.append(Paragraph(metric["rotulo"], styles["Heading3"]))
        story.append(
            Paragraph(
                f"Fórmula: {_text(metric['memoria'].get('formula'))}<br/>"
                f"Componentes: {_text(components)}<br/>"
                f"Fonte: {_source_label(metric['source_refs'])}<br/>"
                f"as_of: {metric.get('as_of') or '—'}",
                styles["SmallAudit"],
            )
        )
        story.append(Spacer(1, 5))

    def _footer(canvas: Any, report: Any) -> None:
        canvas.saveState()
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(colors.HexColor("#63756A"))
        hash_label = document.get("pre_hash", "calculado no fechamento")
        canvas.drawString(15 * mm, 8 * mm, f"Hash do conteúdo: {hash_label}")
        canvas.drawRightString(195 * mm, 8 * mm, f"Página {report.page}")
        canvas.restoreState()

    doc.build(story, onFirstPage=_footer, onLaterPages=_footer)
    return buffer.getvalue()


def _render_xlsx(document: dict[str, Any]) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    summary = workbook.active
    summary.title = "Resumo"
    header = document["cabecalho"]
    summary.append([document["titulo"]])
    summary.append(["Organização", header["organizacao"]])
    summary.append(["Ente", header["ente"]])
    summary.append(["Código IBGE", header["cod_ibge"]])
    summary.append(["Período", document["periodo"]])
    summary.append(["as_of da consulta", document["as_of"]])
    summary.append(["Gerado em", document["gerado_em"]])
    summary.append(["Responsável", header.get("responsavel") or "—"])
    summary["A1"].font = Font(bold=True, size=16, color="173F2A")
    summary.column_dimensions["A"].width = 24
    summary.column_dimensions["B"].width = 72

    metrics = workbook.create_sheet("Indicadores")
    metrics.append(
        [
            "Código",
            "Indicador",
            "Valor",
            "Unidade",
            "Status",
            "Faixa",
            "Fonte",
            "Período fonte",
            "Versão",
            "as_of",
            "Fórmula",
        ]
    )
    for metric in document["metricas"]:
        first_source = metric["source_refs"][0] if metric["source_refs"] else {}
        value = metric.get("valor")
        metrics.append(
            [
                metric["codigo"],
                metric["rotulo"],
                float(value) if value is not None else None,
                metric["unidade"],
                metric["status"],
                metric.get("faixa"),
                _source_label(metric["source_refs"]),
                first_source.get("periodo"),
                first_source.get("versao_entrega"),
                metric.get("as_of"),
                metric["memoria"].get("formula"),
            ]
        )
    _style_sheet(metrics)
    metrics.freeze_panes = "A2"

    issues = workbook.create_sheet("Pendências")
    issues.append(["Tipo", "Código", "Mensagem", "Período esperado", "Período encontrado"])
    for issue in document["dados_incompletos"]:
        issues.append(
            [
                issue["tipo"],
                issue["codigo"],
                issue["mensagem"],
                issue.get("periodo_esperado"),
                issue.get("periodo_encontrado"),
            ]
        )
    _style_sheet(issues, fill="F4D8B5", color="5B3214")

    sources = workbook.create_sheet("Fontes")
    sources.append(["Relatório", "Anexo", "Período", "Versão", "as_of da consulta"])
    for source in document["source_refs"]:
        sources.append(
            [
                source.get("relatorio"),
                source.get("anexo"),
                source.get("periodo"),
                source.get("versao_entrega"),
                document["as_of"],
            ]
        )
    _style_sheet(sources)

    memory = workbook.create_sheet("Memória de cálculo")
    memory.append(["Indicador", "Fórmula", "Componentes", "Fonte", "as_of"])
    for metric in document["metricas"]:
        memory.append(
            [
                metric["rotulo"],
                metric["memoria"].get("formula"),
                json.dumps(metric["memoria"].get("componentes", {}), ensure_ascii=False),
                _source_label(metric["source_refs"]),
                metric.get("as_of"),
            ]
        )
    _style_sheet(memory)

    if document.get("conformidade"):
        compliance = workbook.create_sheet("Conformidade")
        compliance.append(["Relatório", "Período", "Prazo", "Status", "Versão", "Fonte"])
        for item in document["conformidade"]:
            compliance.append(
                [
                    item["relatorio"],
                    item["periodo"],
                    item.get("prazo"),
                    item["status"],
                    item.get("versao_entrega"),
                    _source_label([item["source_ref"]]) if item.get("source_ref") else "—",
                ]
            )
        _style_sheet(compliance)

    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _style_sheet(sheet: Any, *, fill: str = "173F2A", color: str = "FFFFFF") -> None:
    from openpyxl.styles import Alignment, Font, PatternFill

    for cell in sheet[1]:
        cell.fill = PatternFill("solid", fgColor=fill)
        cell.font = Font(bold=True, color=color)
        cell.alignment = Alignment(wrap_text=True, vertical="top")
    for column in sheet.columns:
        letter = column[0].column_letter
        width = min(max(len(_text(cell.value)) for cell in column) + 2, 70)
        sheet.column_dimensions[letter].width = max(width, 12)
        for cell in column:
            cell.alignment = Alignment(wrap_text=True, vertical="top")


def _render_pptx(document: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.core_properties.title = document["titulo"]
    presentation.core_properties.author = (
        document["cabecalho"].get("responsavel") or "Plataforma Fiscal"
    )

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = document["titulo"]
    title_slide.placeholders[1].text = (
        f"{document['cabecalho']['organizacao']}\n{document['cabecalho']['ente']}\n"
        f"Período {document['periodo']} · as_of {document['as_of']} · "
        f"gerado {document['gerado_em']}"
    )

    metrics = document["metricas"]
    for offset in range(0, len(metrics), 6):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Indicadores fiscais"
        chunk = metrics[offset : offset + 6]
        table_shape = slide.shapes.add_table(
            len(chunk) + 1, 5, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.4)
        )
        table = table_shape.table
        headers = ["Indicador", "Valor", "Situação", "Fonte", "as_of"]
        for col, value in enumerate(headers):
            table.cell(0, col).text = value
        for row_idx, metric in enumerate(chunk, start=1):
            values = [
                metric["rotulo"],
                metric["valor_formatado"],
                metric["status"],
                _source_label(metric["source_refs"]),
                metric.get("as_of") or "—",
            ]
            for col, value in enumerate(values):
                table.cell(row_idx, col).text = _text(value)
        _style_ppt_table(table, RGBColor(23, 63, 42), Pt(9))

    audit_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    audit_slide.shapes.title.text = "Rastreabilidade e memória de cálculo"
    body = audit_slide.placeholders[1].text_frame
    body.clear()
    for metric in metrics:
        p = body.add_paragraph()
        p.text = (
            f"{metric['rotulo']}: {metric['memoria'].get('formula', 'valor materializado')} · "
            f"{_source_label(metric['source_refs'])} · as_of {metric.get('as_of') or '—'}"
        )
        p.level = 0
        p.font.size = Pt(12)

    issue_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    issue_slide.shapes.title.text = "Dados incompletos ou defasados"
    issue_body = issue_slide.placeholders[1].text_frame
    issue_body.clear()
    if document["dados_incompletos"]:
        for issue in document["dados_incompletos"]:
            p = issue_body.add_paragraph()
            p.text = f"[{issue['tipo'].upper()}] {issue['codigo']}: {issue['mensagem']}"
            p.font.size = Pt(15)
    else:
        issue_body.text = "Nenhuma incompletude detectada para as seções solicitadas."

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _style_ppt_table(table: Any, header_color: Any, font_size: Any) -> None:
    from pptx.dml.color import RGBColor

    for col in range(len(table.columns)):
        header = table.cell(0, col)
        header.fill.solid()
        header.fill.fore_color.rgb = header_color
        for paragraph in header.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.size = font_size
    for row in range(1, len(table.rows)):
        for col in range(len(table.columns)):
            for paragraph in table.cell(row, col).text_frame.paragraphs:
                paragraph.font.size = font_size
